import pandas as pd
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

# Load analysis results
monthly_stats = pd.read_csv('analysis_results/monthly_stats.csv')
program_click_rates = pd.read_csv('analysis_results/program_click_rates.csv')
program_lpv_rates = pd.read_csv('analysis_results/program_lpv_rates.csv')
program_performance = pd.read_csv('analysis_results/program_performance.csv')
program_roi = pd.read_csv('analysis_results/program_roi.csv')
monthly_trends = pd.read_csv('analysis_results/monthly_trends.csv')

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Define colors
BLUE = RGBColor(0, 51, 102)
GRAY = RGBColor(128, 128, 128)

# ============================================================================
# TITLE SLIDE
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

# Title
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1), Inches(9), Inches(1.5))
title_frame = title_box.text_frame
title_frame.text = "FY2024 Meta Campaign Analysis"
title_para = title_frame.paragraphs[0]
title_para.font.size = Pt(44)
title_para.font.bold = True
title_para.font.color.rgb = BLUE
title_para.alignment = PP_ALIGN.CENTER

# Subtitle
subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(9), Inches(1))
subtitle_frame = subtitle_box.text_frame
subtitle_frame.text = "Lead Generation Campaign Performance Report"
subtitle_para = subtitle_frame.paragraphs[0]
subtitle_para.font.size = Pt(24)
subtitle_para.font.color.rgb = GRAY
subtitle_para.alignment = PP_ALIGN.CENTER

# Date
date_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(9), Inches(0.5))
date_frame = date_box.text_frame
date_frame.text = "December 2024"
date_para = date_frame.paragraphs[0]
date_para.font.size = Pt(18)
date_para.font.color.rgb = GRAY
date_para.alignment = PP_ALIGN.CENTER

# ============================================================================
# EXECUTIVE SUMMARY
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content

title = slide.shapes.title
title.text = "Executive Summary"

content = slide.placeholders[1]
tf = content.text_frame
tf.text = "Key Findings:"
tf.paragraphs[0].font.size = Pt(18)
tf.paragraphs[0].font.bold = True

p = tf.add_paragraph()
p.text = f"• Total Leads Generated: {program_performance['Results'].sum():.0f}"
p.level = 0
p.font.size = Pt(14)

p = tf.add_paragraph()
p.text = f"• Total Spend: ${program_performance['Amount spent (USD)'].sum():,.0f}"
p.level = 0
p.font.size = Pt(14)

p = tf.add_paragraph()
p.text = f"• Overall Click-to-Submit Rate: {program_click_rates['Results'].sum() / program_click_rates['Link clicks'].sum() * 100:.2f}%"
p.level = 0
p.font.size = Pt(14)

p = tf.add_paragraph()
p.text = f"• Overall LPV-to-Submit Rate: {program_lpv_rates['Results'].sum() / program_lpv_rates['Landing page views'].sum() * 100:.2f}%"
p.level = 0
p.font.size = Pt(14)

p = tf.add_paragraph()
p.text = f"• Average CPL: ${program_performance[program_performance['CPL'].notna()]['CPL'].mean():.2f}"
p.level = 0
p.font.size = Pt(14)

p = tf.add_paragraph()
p.text = f"• Programs Analyzed: {len(program_performance)}"
p.level = 0
p.font.size = Pt(14)

# ============================================================================
# QUESTION 1: SEASONALITY ANALYSIS
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[1])

title = slide.shapes.title
title.text = "Seasonality Analysis: CPL and Lead Volume"

content = slide.placeholders[1]
tf = content.text_frame
tf.text = "Key Insights:"
tf.paragraphs[0].font.size = Pt(18)
tf.paragraphs[0].font.bold = True

# Find peak months
peak_leads_month = monthly_stats.loc[monthly_stats['Results'].idxmax(), 'Month_Name']
peak_leads_value = monthly_stats['Results'].max()
low_leads_month = monthly_stats.loc[monthly_stats['Results'].idxmin(), 'Month_Name']
low_leads_value = monthly_stats['Results'].min()

peak_cpl_month = monthly_stats[monthly_stats['CPL'].notna()].loc[monthly_stats[monthly_stats['CPL'].notna()]['CPL'].idxmax(), 'Month_Name']
peak_cpl_value = monthly_stats[monthly_stats['CPL'].notna()]['CPL'].max()
low_cpl_month = monthly_stats[monthly_stats['CPL'].notna()].loc[monthly_stats[monthly_stats['CPL'].notna()]['CPL'].idxmin(), 'Month_Name']
low_cpl_value = monthly_stats[monthly_stats['CPL'].notna()]['CPL'].min()

p = tf.add_paragraph()
p.text = f"• Peak Lead Month: {peak_leads_month} ({peak_leads_value:.0f} leads)"
p.level = 0
p.font.size = Pt(14)

p = tf.add_paragraph()
p.text = f"• Lowest Lead Month: {low_leads_month} ({low_leads_value:.0f} leads)"
p.level = 0
p.font.size = Pt(14)

p = tf.add_paragraph()
p.text = f"• Highest CPL Month: {peak_cpl_month} (${peak_cpl_value:.2f})"
p.level = 0
p.font.size = Pt(14)

p = tf.add_paragraph()
p.text = f"• Lowest CPL Month: {low_cpl_month} (${low_cpl_value:.2f})"
p.level = 0
p.font.size = Pt(14)

p = tf.add_paragraph()
p.text = "\nRecommendation: Increase budget allocation during peak months (May-June) and optimize campaigns during high CPL periods."
p.level = 0
p.font.size = Pt(14)
p.font.italic = True

# Add chart image if exists
if os.path.exists('charts/01_seasonality_analysis.png'):
    slide.shapes.add_picture('charts/01_seasonality_analysis.png', Inches(0.5), Inches(3.5), Inches(9), Inches(3.5))

# ============================================================================
# QUESTION 2: CLICK TO SUBMIT RATE
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[1])

title = slide.shapes.title
title.text = "Click-to-Submit Conversion Rate"

content = slide.placeholders[1]
tf = content.text_frame
total_clicks = program_click_rates['Link clicks'].sum()
total_submits = program_click_rates['Results'].sum()
overall_rate = (total_submits / total_clicks * 100) if total_clicks > 0 else 0

tf.text = f"Overall Rate: {overall_rate:.2f}%"
tf.paragraphs[0].font.size = Pt(20)
tf.paragraphs[0].font.bold = True

p = tf.add_paragraph()
p.text = f"\nTotal Clicks: {total_clicks:,.0f}"
p.level = 0
p.font.size = Pt(14)

p = tf.add_paragraph()
p.text = f"Total Submits: {total_submits:,.0f}"
p.level = 0
p.font.size = Pt(14)

p = tf.add_paragraph()
p.text = "\nTop Performing Programs:"
p.level = 0
p.font.size = Pt(14)
p.font.bold = True

top_programs = program_click_rates.nlargest(3, 'Conversion_Rate')
for idx, row in top_programs.iterrows():
    if row['Results'] > 0:
        p = tf.add_paragraph()
        p.text = f"• {row['Program']}: {row['Conversion_Rate']:.2f}%"
        p.level = 1
        p.font.size = Pt(12)

if os.path.exists('charts/02_click_to_submit_rate.png'):
    slide.shapes.add_picture('charts/02_click_to_submit_rate.png', Inches(0.5), Inches(3.5), Inches(9), Inches(3.5))

# ============================================================================
# QUESTION 3: LPV TO SUBMIT RATE
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[1])

title = slide.shapes.title
title.text = "Landing Page View-to-Submit Conversion Rate"

content = slide.placeholders[1]
tf = content.text_frame
total_lpv = program_lpv_rates['Landing page views'].sum()
total_submits = program_lpv_rates['Results'].sum()
overall_lpv_rate = (total_submits / total_lpv * 100) if total_lpv > 0 else 0

tf.text = f"Overall Rate: {overall_lpv_rate:.2f}%"
tf.paragraphs[0].font.size = Pt(20)
tf.paragraphs[0].font.bold = True

p = tf.add_paragraph()
p.text = f"\nTotal Landing Page Views: {total_lpv:,.0f}"
p.level = 0
p.font.size = Pt(14)

p = tf.add_paragraph()
p.text = f"Total Submits: {total_submits:,.0f}"
p.level = 0
p.font.size = Pt(14)

p = tf.add_paragraph()
p.text = "\nInsight: Landing page optimization shows higher conversion potential than ad clicks alone."
p.level = 0
p.font.size = Pt(14)
p.font.italic = True

if os.path.exists('charts/03_lpv_to_submit_rate.png'):
    slide.shapes.add_picture('charts/03_lpv_to_submit_rate.png', Inches(0.5), Inches(3.5), Inches(9), Inches(3.5))

# ============================================================================
# QUESTION 4: UNDERPERFORMING PROGRAMS
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[1])

title = slide.shapes.title
title.text = "Underperforming Programs Analysis"

content = slide.placeholders[1]
tf = content.text_frame
tf.text = "Programs Requiring Attention:"
tf.paragraphs[0].font.size = Pt(18)
tf.paragraphs[0].font.bold = True

underperforming = program_performance[program_performance['Category'].isin(['Underperforming', 'High CPL', 'Low Conversion', 'No Leads'])]

for idx, row in underperforming.iterrows():
    p = tf.add_paragraph()
    p.text = f"• {row['Program']}"
    p.level = 0
    p.font.size = Pt(14)
    p.font.bold = True
    
    p2 = tf.add_paragraph()
    p2.text = f"  Category: {row['Category']} | CPL: ${row['CPL']:.2f} | Conversion: {row['Click_to_Submit_Rate']:.2f}%"
    p2.level = 1
    p2.font.size = Pt(12)

p = tf.add_paragraph()
p.text = "\nRecommendation: Review targeting, ad creative, and landing page experience for underperforming programs."
p.level = 0
p.font.size = Pt(14)
p.font.italic = True

if os.path.exists('charts/04_underperforming_programs.png'):
    slide.shapes.add_picture('charts/04_underperforming_programs.png', Inches(0.5), Inches(4), Inches(9), Inches(3))

# ============================================================================
# QUESTION 5: TRENDS ANALYSIS
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[1])

title = slide.shapes.title
title.text = "Ad Engagement Trends & Lead Correlation"

content = slide.placeholders[1]
tf = content.text_frame
tf.text = "Key Findings:"
tf.paragraphs[0].font.size = Pt(18)
tf.paragraphs[0].font.bold = True

p = tf.add_paragraph()
p.text = "• Strong positive correlation between ad engagements and leads in following periods"
p.level = 0
p.font.size = Pt(14)

p = tf.add_paragraph()
p.text = "• Clicks (1-week lag) correlation: 0.651"
p.level = 0
p.font.size = Pt(14)

p = tf.add_paragraph()
p.text = "• Landing Page Views (1-week lag) correlation: 0.661"
p.level = 0
p.font.size = Pt(14)

p = tf.add_paragraph()
p.text = "\nRecommendation: Increase ad spend during high-engagement periods to capitalize on delayed conversion patterns."
p.level = 0
p.font.size = Pt(14)
p.font.italic = True

if os.path.exists('charts/05_engagement_trends.png'):
    slide.shapes.add_picture('charts/05_engagement_trends.png', Inches(0.5), Inches(3.5), Inches(9), Inches(3.5))

# ============================================================================
# QUESTION 6: BUDGET ALLOCATION RECOMMENDATIONS
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[1])

title = slide.shapes.title
title.text = "Budget Allocation Recommendations"

content = slide.placeholders[1]
tf = content.text_frame
tf.text = "Recommended Budget Reallocation:"
tf.paragraphs[0].font.size = Pt(18)
tf.paragraphs[0].font.bold = True

# Get top performers
top_performers = program_roi[program_roi['Efficiency_Score'].notna()].nlargest(3, 'Efficiency_Score')

p = tf.add_paragraph()
p.text = "\nINCREASE Budget:"
p.level = 0
p.font.size = Pt(14)
p.font.bold = True

for idx, row in top_performers.iterrows():
    p = tf.add_paragraph()
    p.text = f"• {row['Program']} (Efficiency Score: {row['Efficiency_Score']:.3f})"
    p.level = 1
    p.font.size = Pt(12)

# Get underperformers
underperformers = program_roi[program_roi['Category'].isin(['High CPL', 'No Leads', 'Low Conversion'])].nsmallest(2, 'Efficiency_Score')

p = tf.add_paragraph()
p.text = "\nDECREASE/OPTIMIZE Budget:"
p.level = 0
p.font.size = Pt(14)
p.font.bold = True

for idx, row in underperformers.iterrows():
    p = tf.add_paragraph()
    p.text = f"• {row['Program']} ({row['Category']})"
    p.level = 1
    p.font.size = Pt(12)

if os.path.exists('charts/06_budget_allocation.png'):
    slide.shapes.add_picture('charts/06_budget_allocation.png', Inches(0.5), Inches(4.5), Inches(9), Inches(2.5))

# ============================================================================
# QUESTION 7: FY2025 RECOMMENDATIONS
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[1])

title = slide.shapes.title
title.text = "FY2025 Strategic Recommendations"

content = slide.placeholders[1]
tf = content.text_frame
tf.text = "Key Recommendations:"
tf.paragraphs[0].font.size = Pt(18)
tf.paragraphs[0].font.bold = True

p = tf.add_paragraph()
p.text = "1. Seasonal Budget Planning"
p.level = 0
p.font.size = Pt(14)
p.font.bold = True

p2 = tf.add_paragraph()
p2.text = "   Allocate 40-50% of Q4 budget to May-June period (peak lead months)"
p2.level = 1
p2.font.size = Pt(12)

p = tf.add_paragraph()
p.text = "2. Program Optimization"
p.level = 0
p.font.size = Pt(14)
p.font.bold = True

p2 = tf.add_paragraph()
p2.text = "   Pause or significantly reduce spend on 'Masters of Space Operations' (0 leads)"
p2.level = 1
p2.font.size = Pt(12)

p2 = tf.add_paragraph()
p2.text = "   Optimize 'MS in Engineering Management' (highest CPL at $177.39)"
p2.level = 1
p2.font.size = Pt(12)

p = tf.add_paragraph()
p.text = "3. Landing Page Optimization"
p.level = 0
p.font.size = Pt(14)
p.font.bold = True

p2 = tf.add_paragraph()
p2.text = "   Focus on improving LPV-to-submit rate (currently 2.81%)"
p2.level = 1
p2.font.size = Pt(12)

p = tf.add_paragraph()
p.text = "4. Engagement-Based Budgeting"
p.level = 0
p.font.size = Pt(14)
p.font.bold = True

p2 = tf.add_paragraph()
p2.text = "   Increase spend during high engagement weeks to capture delayed conversions"
p2.level = 1
p2.font.size = Pt(12)

# ============================================================================
# CONCLUSION SLIDE
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[1])

title = slide.shapes.title
title.text = "Conclusion & Next Steps"

content = slide.placeholders[1]
tf = content.text_frame
tf.text = "Summary:"
tf.paragraphs[0].font.size = Pt(18)
tf.paragraphs[0].font.bold = True

p = tf.add_paragraph()
p.text = "• Strong seasonality patterns identified - peak performance in May-June"
p.level = 0
p.font.size = Pt(14)

p = tf.add_paragraph()
p.text = "• Conversion rates show room for improvement, especially landing page optimization"
p.level = 0
p.font.size = Pt(14)

p = tf.add_paragraph()
p.text = "• Clear opportunities for budget reallocation to high-performing programs"
p.level = 0
p.font.size = Pt(14)

p = tf.add_paragraph()
p.text = "• Engagement trends indicate delayed conversion patterns to leverage"
p.level = 0
p.font.size = Pt(14)

p = tf.add_paragraph()
p.text = "\nNext Steps:"
p.level = 0
p.font.size = Pt(14)
p.font.bold = True

p = tf.add_paragraph()
p.text = "1. Implement seasonal budget allocation plan"
p.level = 1
p.font.size = Pt(12)

p = tf.add_paragraph()
p.text = "2. Conduct landing page A/B testing for underperforming programs"
p.level = 1
p.font.size = Pt(12)

p = tf.add_paragraph()
p.text = "3. Review and optimize 'Masters of Space Operations' campaign strategy"
p.level = 1
p.font.size = Pt(12)

p = tf.add_paragraph()
p.text = "4. Set up engagement-based budget triggers for Q1 2025"
p.level = 1
p.font.size = Pt(12)

# Save presentation
prs.save('FY2024_Campaign_Analysis_Presentation.pptx')
print("Presentation saved as 'FY2024_Campaign_Analysis_Presentation.pptx'")

