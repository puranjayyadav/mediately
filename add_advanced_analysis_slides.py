import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

# Load existing presentation
prs = Presentation('FY2024_Campaign_Analysis_Presentation.pptx')

# Define colors
BLUE = RGBColor(0, 51, 102)
GRAY = RGBColor(128, 128, 128)

# Load predictive analysis results
roi_scenarios = pd.read_csv('analysis_results/roi_optimization.csv')
what_if = pd.read_csv('analysis_results/what_if_scenarios.csv')
program_roi = pd.read_csv('analysis_results/program_roi.csv')

# ============================================================================
# PREDICTIVE ANALYSIS SLIDE
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[1])

title = slide.shapes.title
title.text = "Predictive Analysis: FY2025 Forecast"

content = slide.placeholders[1]
tf = content.text_frame
tf.text = "6-Month Forecast (Jul-Dec 2025):"
tf.paragraphs[0].font.size = Pt(18)
tf.paragraphs[0].font.bold = True

p = tf.add_paragraph()
p.text = "• December 2025: Best opportunity (44 leads, $98.06 CPL)"
p.level = 0
p.font.size = Pt(14)

p = tf.add_paragraph()
p.text = "• September & November: Low volume, high CPL expected"
p.level = 0
p.font.size = Pt(14)

p = tf.add_paragraph()
p.text = "• Q4 (Oct-Dec) forecasted stronger than Q1 (Jul-Sep)"
p.level = 0
p.font.size = Pt(14)

p = tf.add_paragraph()
p.text = "\nRecommendation: Increase December budget by 20-30% to capitalize on forecasted strong performance."
p.level = 0
p.font.size = Pt(14)
p.font.italic = True

if os.path.exists('charts/07_predictive_forecast.png'):
    slide.shapes.add_picture('charts/07_predictive_forecast.png', Inches(0.5), Inches(3.5), Inches(9), Inches(3.5))

# ============================================================================
# OPTIMAL BUDGET ALLOCATION SLIDE
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[1])

title = slide.shapes.title
title.text = "Prescriptive Analysis: Optimal Budget Allocation"

content = slide.placeholders[1]
tf = content.text_frame
tf.text = "Balanced Approach (Recommended):"
tf.paragraphs[0].font.size = Pt(18)
tf.paragraphs[0].font.bold = True

p = tf.add_paragraph()
p.text = "• Predicted Leads: 1,890 (vs. current 619)"
p.level = 0
p.font.size = Pt(14)

p = tf.add_paragraph()
p.text = "• Improvement: +205% lead volume"
p.level = 0
p.font.size = Pt(14)

p = tf.add_paragraph()
p.text = "• Key Changes:"
p.level = 0
p.font.size = Pt(14)
p.font.bold = True

p = tf.add_paragraph()
p.text = "  - Increase MS in Management: $1,500 → $28,123"
p.level = 1
p.font.size = Pt(12)

p = tf.add_paragraph()
p.text = "  - Increase Military Campaign: Maintain high investment"
p.level = 1
p.font.size = Pt(12)

p = tf.add_paragraph()
p.text = "  - Reduce MSEM: $26,086 → $18,897 (optimize first)"
p.level = 1
p.font.size = Pt(12)

if os.path.exists('charts/08_optimal_budget_allocation.png'):
    slide.shapes.add_picture('charts/08_optimal_budget_allocation.png', Inches(0.5), Inches(4.5), Inches(9), Inches(2.5))

# ============================================================================
# WHAT-IF SCENARIOS SLIDE
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[1])

title = slide.shapes.title
title.text = "What-If Scenario Analysis"

content = slide.placeholders[1]
tf = content.text_frame
tf.text = "Scenario Comparison:"
tf.paragraphs[0].font.size = Pt(18)
tf.paragraphs[0].font.bold = True

best_scenario = what_if.loc[what_if['Total_Leads'].idxmax()]

p = tf.add_paragraph()
p.text = f"Best Scenario: {best_scenario['Scenario']}"
p.level = 0
p.font.size = Pt(14)
p.font.bold = True

p = tf.add_paragraph()
p.text = f"• Predicted Leads: {best_scenario['Total_Leads']:.0f} (+{best_scenario['Improvement_Leads']:.1f}%)"
p.level = 1
p.font.size = Pt(12)

p = tf.add_paragraph()
p.text = f"• Average CPL: ${best_scenario['Avg_CPL']:.2f} ({best_scenario['Improvement_CPL']:.1f}% improvement)"
p.level = 1
p.font.size = Pt(12)

p = tf.add_paragraph()
p.text = f"• Budget: ${best_scenario['Total_Budget']:,.0f}"
p.level = 1
p.font.size = Pt(12)

p = tf.add_paragraph()
p.text = "\nRecommendation: Implement 'Pause Underperformers, Reallocate' scenario for maximum impact."
p.level = 0
p.font.size = Pt(14)
p.font.italic = True

if os.path.exists('charts/09_what_if_scenarios.png'):
    slide.shapes.add_picture('charts/09_what_if_scenarios.png', Inches(0.5), Inches(4), Inches(9), Inches(3))

# ============================================================================
# ROI OPTIMIZATION SLIDE
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[1])

title = slide.shapes.title
title.text = "ROI Optimization Model"

content = slide.placeholders[1]
tf = content.text_frame
tf.text = "Budget Scaling Analysis:"
tf.paragraphs[0].font.size = Pt(18)
tf.paragraphs[0].font.bold = True

optimal_budget = roi_scenarios[roi_scenarios['Budget_Multiplier'] == 1.0].iloc[0]

p = tf.add_paragraph()
p.text = f"• Optimal Budget Range: ${roi_scenarios['Total_Budget'].min():,.0f} - ${roi_scenarios['Total_Budget'].max():,.0f}"
p.level = 0
p.font.size = Pt(14)

p = tf.add_paragraph()
p.text = f"• ROI: {optimal_budget['ROI']:.4f} (constant across budget levels)"
p.level = 0
p.font.size = Pt(14)

p = tf.add_paragraph()
p.text = f"• Predicted Leads: {optimal_budget['Total_Leads']:.0f} - {roi_scenarios['Total_Leads'].max():.0f}"
p.level = 0
p.font.size = Pt(14)

p = tf.add_paragraph()
p.text = "\nInsight: Budget scales linearly with leads - each 10% increase yields ~189 additional leads."
p.level = 0
p.font.size = Pt(14)
p.font.italic = True

if os.path.exists('charts/10_roi_optimization.png'):
    slide.shapes.add_picture('charts/10_roi_optimization.png', Inches(0.5), Inches(3.5), Inches(9), Inches(3.5))

# ============================================================================
# UPDATED STRATEGIC RECOMMENDATIONS SLIDE
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[1])

title = slide.shapes.title
title.text = "Enhanced FY2025 Strategic Recommendations"

content = slide.placeholders[1]
tf = content.text_frame
tf.text = "Based on Predictive & Prescriptive Analysis:"
tf.paragraphs[0].font.size = Pt(18)
tf.paragraphs[0].font.bold = True

p = tf.add_paragraph()
p.text = "1. Optimal Budget Allocation"
p.level = 0
p.font.size = Pt(14)
p.font.bold = True

p2 = tf.add_paragraph()
p2.text = "   Implement balanced approach: $162,531 budget → 1,890 leads (+205%)"
p2.level = 1
p2.font.size = Pt(12)

p = tf.add_paragraph()
p.text = "2. What-If Scenario Implementation"
p.level = 0
p.font.size = Pt(14)
p.font.bold = True

p2 = tf.add_paragraph()
p2.text = "   Pause underperformers, reallocate to top programs"
p2.level = 1
p2.font.size = Pt(12)

p2 = tf.add_paragraph()
p2.text = "   Expected: 1,934 leads, $84.02 CPL (+212.5% improvement)"
p2.level = 1
p2.font.size = Pt(12)

p = tf.add_paragraph()
p.text = "3. Predictive Budget Planning"
p.level = 0
p.font.size = Pt(14)
p.font.bold = True

p2 = tf.add_paragraph()
p2.text = "   Increase December 2025 budget by 20-30% (forecasted best month)"
p2.level = 1
p2.font.size = Pt(12)

p = tf.add_paragraph()
p.text = "4. ROI Optimization"
p.level = 0
p.font.size = Pt(14)
p.font.bold = True

p2 = tf.add_paragraph()
p2.text = "   Maintain budget range: $162,531 - $178,784 for optimal ROI"
p2.level = 1
p2.font.size = Pt(12)

# ============================================================================
# EXPECTED IMPACT SLIDE
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[1])

title = slide.shapes.title
title.text = "Expected Impact Summary"

content = slide.placeholders[1]
tf = content.text_frame
tf.text = "If All Recommendations Are Implemented:"
tf.paragraphs[0].font.size = Pt(18)
tf.paragraphs[0].font.bold = True

p = tf.add_paragraph()
p.text = "• Lead Volume: 619 → 1,890-2,079 (+205-236%)"
p.level = 0
p.font.size = Pt(14)

p = tf.add_paragraph()
p.text = "• Average CPL: $95.77 → $84.02-$85.99 (-10-12%)"
p.level = 0
p.font.size = Pt(14)

p = tf.add_paragraph()
p.text = "• Total Budget: $209,548 → $162,531-$178,784 (-15-22%)"
p.level = 0
p.font.size = Pt(14)

p = tf.add_paragraph()
p.text = "• ROI: 0.0030 → 0.0116 (+287%)"
p.level = 0
p.font.size = Pt(14)

p = tf.add_paragraph()
p.text = "\nKey Benefits:"
p.level = 0
p.font.size = Pt(14)
p.font.bold = True

p = tf.add_paragraph()
p.text = "✓ 3x increase in lead volume"
p.level = 1
p.font.size = Pt(12)

p = tf.add_paragraph()
p.text = "✓ 10-12% reduction in CPL"
p.level = 1
p.font.size = Pt(12)

p = tf.add_paragraph()
p.text = "✓ 15-22% budget reduction while improving outcomes"
p.level = 1
p.font.size = Pt(12)

p = tf.add_paragraph()
p.text = "✓ Nearly 4x improvement in ROI"
p.level = 1
p.font.size = Pt(12)

# Save updated presentation
prs.save('FY2024_Campaign_Analysis_Presentation_Enhanced.pptx')
print("Enhanced presentation saved as 'FY2024_Campaign_Analysis_Presentation_Enhanced.pptx'")

